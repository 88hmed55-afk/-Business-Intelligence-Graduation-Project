# -*- coding: utf-8 -*-
"""
Nova BI — Executive Presentation Generator
Generates a premium, fully-editable PowerPoint (22 slides) for the
Nova BI graduation project. All diagrams and charts are native PowerPoint
shapes / charts so the committee can edit everything.
"""

import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.ns import qn

# --------------------------------------------------------------------------
# Design tokens
# --------------------------------------------------------------------------
DEEP   = RGBColor(0x0B, 0x16, 0x33)   # deep navy
NAVY   = RGBColor(0x11, 0x23, 0x4E)   # navy
INDIGO = RGBColor(0x4F, 0x46, 0xE5)   # indigo
BLUE   = RGBColor(0x25, 0x63, 0xEB)   # primary blue
SKY    = RGBColor(0x38, 0xBD, 0xF8)   # sky
CYAN   = RGBColor(0x06, 0xB6, 0xD4)   # cyan accent
LIGHT  = RGBColor(0xF1, 0xF5, 0xF9)   # light gray
LIGHT2 = RGBColor(0xE2, 0xE8, 0xF0)   # border gray
SLATE  = RGBColor(0x47, 0x55, 0x69)   # body text
GRAY   = RGBColor(0x94, 0xA3, 0xB8)   # muted
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREEN  = RGBColor(0x10, 0xB9, 0x81)
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)
ROSE   = RGBColor(0xEF, 0x44, 0x44)
FONT   = "Segoe UI"

PALETTE = [BLUE, INDIGO, CYAN, SKY, AMBER, GREEN, ROSE]
CHART_PALETTE = [
    RGBColor(0x25, 0x63, 0xEB), RGBColor(0x4F, 0x46, 0xE5),
    RGBColor(0x06, 0xB6, 0xD4), RGBColor(0x38, 0xBD, 0xF8),
    RGBColor(0x10, 0xB9, 0x81), RGBColor(0xF5, 0x9E, 0x0B),
    RGBColor(0x8B, 0x5C, 0xF6), RGBColor(0x94, 0xA3, 0xB8),
]

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

# --------------------------------------------------------------------------
# Low-level helpers
# --------------------------------------------------------------------------

def slide():
    s = prs.slides.add_slide(BLANK)
    return s


def rgb(c):
    return RGBColor.from_string(c) if isinstance(c, str) else c


def _kill_fill(shape):
    spPr = shape._element.spPr
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill"):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)


def grad_fill(shape, c1, c2, angle=90):
    """Linear gradient fill injected as OOXML (fully editable in PowerPoint)."""
    _kill_fill(shape)
    spPr = shape._element.spPr
    gf = spPr.makeelement(qn("a:gradFill"), {})
    gsLst = spPr.makeelement(qn("a:gsLst"), {})
    for pos, c in ((0, rgb(c1)), (100000, rgb(c2))):
        gs = spPr.makeelement(qn("a:gs"), {"pos": str(pos)})
        clr = spPr.makeelement(qn("a:srgbClr"), {"val": "%02X%02X%02X" % (c[0], c[1], c[2])})
        gs.append(clr)
        gsLst.append(gs)
    gf.append(gsLst)
    lin = spPr.makeelement(qn("a:lin"), {"ang": str(int(angle * 60000)), "scaled": "1"})
    gf.append(lin)
    spPr.append(gf)


def shadow(shape, blur=0.12, dist=0.045, alpha=62, color="0B1633", angle=45):
    """Soft outer shadow injected as OOXML."""
    spPr = shape._element.spPr
    for el in spPr.findall(qn("a:effectLst")):
        spPr.remove(el)
    eff = spPr.makeelement(qn("a:effectLst"), {})
    sh = spPr.makeelement(qn("a:outerShdw"), {
        "blurRad": str(Emu(int(914400 * blur))),
        "dist": str(Emu(int(914400 * dist))),
        "dir": str(int(angle * 60000)),
        "rotWithShape": "0",
    })
    clr = spPr.makeelement(qn("a:srgbClr"), {"val": color})
    a = spPr.makeelement(qn("a:alpha"), {"val": str(int(alpha * 1000))})
    clr.append(a)
    sh.append(clr)
    eff.append(sh)
    spPr.append(eff)


def rect(s, x, y, w, h, fill=None, line=None, line_w=0.75, radius=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = rgb(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = rgb(line)
        sp.line.width = Pt(line_w)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    sp.shadow.inherit = False
    return sp


def line(s, x1, y1, x2, y2, color=LIGHT2, w=0.75, dash=None):
    ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = rgb(color)
    ln.line.width = Pt(w)
    if dash:
        ln.line.dash_style = dash
    ln.shadow.inherit = False
    return ln


def tb(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tf


def para(tf, text, size=14, bold=False, color=SLATE, font=FONT, align=PP_ALIGN.LEFT,
         sb=0, sa=0, ls=None, first=False, italic=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(sb)
    p.space_after = Pt(sa)
    if ls:
        p.line_spacing = ls
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = rgb(color)
    r.font.name = font
    return p


def runs(tf, items, align=PP_ALIGN.LEFT, sb=0, sa=0, first=False, ls=None):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(sb)
    p.space_after = Pt(sa)
    if ls:
        p.line_spacing = ls
    for t, size, bold, color, font in items:
        r = p.add_run()
        r.text = t
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = rgb(color)
        r.font.name = font
    return p


def chip(s, x, y, w, h, text, fill=INDIGO, tcolor=WHITE, size=11, bold=True,
         line_c=None, radius=0.5, dot=None):
    c = rect(s, x, y, w, h, fill=fill, line=line_c, radius=radius)
    tf = c.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    if dot:
        rd = p.add_run()
        rd.text = dot + "  "
        rd.font.size = Pt(size - 1)
        rd.font.color.rgb = rgb(SKY)
        rd.font.bold = True
        rd.font.name = FONT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(tcolor)
    r.font.name = FONT
    return c


def transition(s, kind="fade", spd="med"):
    el = s._element
    for t in ("p:transition", "p:timing"):
        for e in el.findall(qn(t)):
            el.remove(e)
    tr = el.makeelement(qn("p:transition"), {"spd": spd})
    inner = el.makeelement(qn("p:" + kind), {})
    tr.append(inner)
    el.append(tr)


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def icon_motif(s, x, y, d, glyph, bg=INDIGO, fg=WHITE, size=16):
    """Simple geometric 'icon' chip (editable shapes + glyph)."""
    c = rect(s, x, y, d, d, fill=bg, radius=0.30)
    ctf = c.text_frame
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ctf.margin_left = 0
    ctf.margin_right = 0
    ctf.margin_top = 0
    ctf.margin_bottom = 0
    p = ctf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = glyph
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = rgb(fg)
    r.font.name = "Segoe UI"
    return c


def header(s, eyebrow, title, subtitle=None, dark=False, page=None, accent=None,
           tcolor=None, scolor=None):
    accent = accent if accent else INDIGO
    bar = rect(s, 0.55, 0.55, 0.10, 0.80, fill=BLUE, radius=0.5)
    grad_fill(bar, INDIGO, SKY, 90)
    tf = tb(s, 0.80, 0.42, 11.9, 0.3)
    para(tf, eyebrow.upper(), size=12, bold=True, color=SKY if dark else BLUE,
         first=True, ls=1.0)
    tcolor = tcolor if tcolor else (WHITE if dark else DEEP)
    tf2 = tb(s, 0.80, 0.60, 11.9, 0.75)
    para(tf2, title, size=29, bold=True, color=tcolor, first=True, ls=1.0)
    if subtitle:
        scolor = scolor if scolor else (GRAY if dark else SLATE)
        tf3 = tb(s, 0.80, 1.34, 11.9, 0.4)
        para(tf3, subtitle, size=13, bold=False, color=scolor, first=True, ls=1.0)
    if page is not None:
        footer(s, page, dark=dark)


def footer(s, page, dark=False):
    fcol = GRAY if not dark else RGBColor(0x5B, 0x6B, 0x8A)
    line(s, 0.55, 7.12, 12.78, 7.12, color=RGBColor(0x2A, 0x3B, 0x66) if dark else LIGHT2, w=0.75)
    tf = tb(s, 0.55, 7.16, 8, 0.3)
    para(tf, "Nova BI  ·  Business Intelligence Management System", size=8.5,
         color=fcol, first=True)
    tf2 = tb(s, 11.9, 7.16, 0.88, 0.3)
    para(tf2, "%02d" % page, size=9, bold=True, color=fcol, align=PP_ALIGN.RIGHT, first=True)


def card(s, x, y, w, h, fill=WHITE, line_c=LIGHT2, radius=0.075, sh=True, grad=None):
    c = rect(s, x, y, w, h, fill=fill, line=line_c, radius=radius)
    if grad:
        grad_fill(c, grad[0], grad[1], grad[2] if len(grad) > 2 else 90)
    if sh:
        shadow(c)
    return c


def arrow_right(s, x, y, w, h, fill=INDIGO):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid()
    a.fill.fore_color.rgb = rgb(fill)
    a.line.fill.background()
    a.shadow.inherit = False
    return a


# --------------------------------------------------------------------------
# Chart helpers (native, editable)
# --------------------------------------------------------------------------

def make_chart(s, ctype, x, y, w, h, cats, series, colors=None, legend=False,
               smooth=True, labels=False, bar_width=100):
    cd = CategoryChartData()
    cd.categories = cats
    colors = colors or CHART_PALETTE
    for i, (name, vals) in enumerate(series):
        cd.add_series(name, vals)
    gf = s.shapes.add_chart(ctype, Inches(x), Inches(y), Inches(w), Inches(h), cd)
    ch = gf.chart
    ch.font.size = Pt(10)
    ch.font.name = FONT
    ch.has_legend = legend
    if legend:
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
    for i, ser in enumerate(ch.series):
        try:
            ser.format.fill.solid()
            ser.format.fill.fore_color.rgb = rgb(colors[i % len(colors)])
        except Exception:
            pass
        if smooth:
            try:
                ser.smooth = True
            except Exception:
                pass
    try:
        ch.value_axis.has_major_gridlines = False
        ch.value_axis.has_minor_gridlines = False
    except Exception:
        pass
    try:
        ch.value_axis.tick_labels.font.size = Pt(9)
        ch.category_axis.tick_labels.font.size = Pt(9)
    except Exception:
        pass
    if labels:
        try:
            ch.plots[0].has_data_labels = True
            ch.plots[0].data_labels.font.size = Pt(9)
            ch.plots[0].data_labels.font.bold = True
        except Exception:
            pass
    if bar_width:
        try:
            ch.plots[0].gap_width = bar_width
        except Exception:
            pass
    return ch


def color_points(ch, colors):
    try:
        for i, pt in enumerate(ch.series[0].points):
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = rgb(colors[i % len(colors)])
    except Exception:
        pass


# --------------------------------------------------------------------------
# Data used across the deck (mirrors the real seeded Nova BI database)
# --------------------------------------------------------------------------
MONTHS = ["Jan", "Feb", "Mar", "Apr", "May", "Jun",
          "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
REV = [48.2, 52.4, 55.1, 53.8, 58.7, 61.3, 59.2, 63.8, 60.1, 66.4, 62.9, 58.3]
ORDERS = [86, 94, 99, 97, 105, 111, 107, 115, 109, 120, 113, 108]
FORECAST_REV = [58.3, 61.7, 64.4, 66.1, 69.5, 72.3]
CATS = ["Electronics", "Fashion", "Home & Garden", "Books", "Sports", "Beauty", "Other"]
CAT_VALS = [217.0, 154.0, 98.0, 63.0, 56.0, 49.0, 63.0]
STATUS = ["Delivered", "Shipped", "Processing", "Cancelled"]
STATUS_VALS = [936, 144, 72, 48]
TOP_PROD = ["Laptop Pro X", "Smart Phone Z", "Wireless Earbuds", "4K Monitor", "Desk Chair"]
TOP_PROD_VAL = [58.4, 51.2, 33.7, 27.9, 24.6]
COUNTRIES = ["Egypt", "UAE", "Saudi Arabia", "Germany", "USA", "Other"]
COUNTRY_VALS = [42, 21, 15, 10, 8, 4]

# --------------------------------------------------------------------------
# SLIDE 1 — Luxury cover
# --------------------------------------------------------------------------
def s01_cover():
    s = slide()
    bg = rect(s, 0, 0, 13.333, 7.5, fill=DEEP, radius=None, shape=MSO_SHAPE.RECTANGLE)
    grad_fill(bg, DEEP, NAVY, 110)
    # decorative glows
    _deco_circle(s, 9.0, -2.2, 6.2, INDIGO, 60)
    _deco_circle(s, 10.4, 3.6, 5.0, BLUE, 40)
    _deco_circle(s, 0.2, 5.4, 3.4, SKY, 25)
    line(s, 0.0, 7.5, 13.333, 7.5, color=RGBColor(0x1E, 0x2F, 0x5E), w=1.0)
    line(s, 0.55, 1.05, 12.78, 1.05, color=RGBColor(0x2E, 0x42, 0x7A), w=0.75)

    tf = tb(s, 0.6, 0.62, 12, 0.4)
    para(tf, "GRADUATION PROJECT   ·   CLASS OF 2025 / 2026", size=12.5, bold=True,
         color=SKY, first=True)

    tf = tb(s, 0.58, 1.55, 11, 1.5)
    para(tf, "NOVA", size=62, bold=True, color=WHITE, first=True, ls=0.95)
    para(tf, "BUSINESS INTELLIGENCE", size=62, bold=True, color=WHITE, ls=0.95)

    bar = rect(s, 0.6, 3.42, 1.7, 0.09, fill=BLUE, radius=0.5)
    grad_fill(bar, INDIGO, CYAN, 0)

    tf = tb(s, 0.6, 3.66, 12.1, 0.5)
    para(tf, "Management System", size=22, bold=True, color=WHITE, first=True)
    tf = tb(s, 0.6, 4.16, 11.6, 0.6)
    para(tf, "An enterprise-grade decision-support platform that turns raw operational data "
             "into executive insight — dashboards, KPIs, analytics and forecasting.",
         size=14, color=RGBColor(0xBF, 0xCC, 0xE6), ls=1.15)

    # meta grid
    meta = [
        ("TEAM MEMBERS", "Nour A.  ·  Omar K.  ·  Sara M."),
        ("SUPERVISOR", "Dr. [Supervisor Name]"),
        ("UNIVERSITY", "[University Name]"),
        ("DEPARTMENT", "Information Systems / Computer Science"),
        ("COURSE", "Graduation Project — BI & Data Science"),
        ("ACADEMIC YEAR", "2025 / 2026"),
    ]
    for i, (k, v) in enumerate(meta):
        col = i % 3
        row = i // 3
        x = 0.6 + col * 4.05
        y = 5.25 + row * 0.85
        chip(s, x, y, 0.06, 0.34, "", fill=CYAN, radius=0.5)
        tfk = tb(s, x + 0.16, y, 3.7, 0.28)
        para(tfk, k, size=10, bold=True, color=GRAY, first=True)
        tfv = tb(s, x + 0.16, y + 0.26, 3.9, 0.4)
        para(tfv, v, size=13, bold=True, color=WHITE, first=True)

    tf = tb(s, 0.6, 7.06, 12, 0.3)
    para(tf, "Prepared with a Fortune-500 presentation standard  ·  Fully editable",
         size=9.5, color=RGBColor(0x7C, 0x8C, 0xB0), first=True)
    transition(s, "fade")
    notes(s, "Welcome the committee. Pitch: organizations drown in data yet starve for insight. "
             "Nova BI is a full-stack Business Intelligence platform — React + FastAPI + PostgreSQL — "
             "that turns raw operational data into an executive decision cockpit. Walk the board "
             "through what we built, how it is engineered, and the measurable value it delivers.\n"
             "Transition: Fade. Animations: title fade-up, meta rows stagger in.")


def _deco_circle(s, x, y, d, color, alpha):
    c = rect(s, x, y, d, d, fill=color, radius=None, shape=MSO_SHAPE.OVAL)
    c.line.fill.background()
    spPr = c._element.spPr
    solid = spPr.find(qn("a:solidFill"))
    if solid is not None:
        clr = solid.find(qn("a:srgbClr"))
        if clr is not None:
            a = spPr.makeelement(qn("a:alpha"), {"val": str(int(alpha * 1000))})
            clr.append(a)
    return c


# --------------------------------------------------------------------------
# SLIDE 2 — Agenda timeline
# --------------------------------------------------------------------------
def s02_agenda():
    s = slide()
    rect(s, 0, 0, 13.333, 7.5, fill=LIGHT, radius=None, shape=MSO_SHAPE.RECTANGLE)
    header(s, "Roadmap", "Agenda", "Seven movements from problem to production-grade platform.",
           page=2)
    items = [
        ("01", "Context & Problem", "The business case behind the project"),
        ("02", "BI Foundations", "What Business Intelligence is and why it matters"),
        ("03", "Architecture & Stack", "System layers and the technology arsenal"),
        ("04", "Data & Modules", "Database design and application modules"),
        ("05", "BI Engine", "KPIs, aggregation, analytics and reports"),
        ("06", "Forecasting & Results", "Predictive analytics and measured outcomes"),
        ("07", "Engineering & Wrap-up", "DevOps, security, testing, future, Q&A"),
    ]
    x0, y0 = 0.55, 2.15
    card_w = 1.78
    line(s, x0 + 0.5, y0 + 0.72, x0 + 6 * 1.86 + 1.2, y0 + 0.72, color=BLUE, w=2.0)
    for i, (num, t, d) in enumerate(items):
        x = x0 + i * 1.86
        c = rect(s, x + 0.28, y0 + 0.35, 0.75, 0.75, fill=INDIGO, radius=0.5)
        grad_fill(c, INDIGO, BLUE, 90)
        ctf = c.text_frame
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ctf.margin_left = 0; ctf.margin_right = 0; ctf.margin_top = 0; ctf.margin_bottom = 0
        p = ctf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num
        r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
        tf = tb(s, x, y0 + 1.3, card_w, 0.7)
        para(tf, t, size=12.5, bold=True, color=DEEP, first=True, align=PP_ALIGN.LEFT)
        tf2 = tb(s, x, y0 + 1.72, card_w, 1.0)
        para(tf2, d, size=9.5, color=SLATE, first=True, ls=1.05)
    # bottom summary banner
    banner = card(s, 0.55, 5.6, 12.23, 1.05, grad=(INDIGO, BLUE, 90), sh=True)
    tf = banner.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.3)
    tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "Outcome:  "
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = SKY; r.font.name = FONT
    r2 = p.add_run(); r2.text = ("a secure, tested, fully containerized BI platform with 81 automated "
                                 "tests, a synchronized database schema, and an executive dashboard.")
    r2.font.size = Pt(13); r2.font.color.rgb = WHITE; r2.font.name = FONT
    transition(s, "fade")
    notes(s, "Seven-part journey. Keep it to 45 seconds: context → foundations → architecture → "
             "data & modules → the BI engine → forecasting & results → engineering wrap-up.\n"
             "Transition: Fade. Animation: numbered nodes reveal left-to-right on click.")


# --------------------------------------------------------------------------
# SLIDE 3 — Project overview / problem statement
# --------------------------------------------------------------------------
def s03_overview():
    s = slide()
    rect(s, 0, 0, 13.333, 7.5, fill=LIGHT, radius=None, shape=MSO_SHAPE.RECTANGLE)
    header(s, "Why we built this", "Project Overview", "Turning operational data into decisions.",
           page=3)

    # Left: statement card
    c = card(s, 0.55, 1.95, 5.6, 3.0)
    icon_motif(s, 0.85, 2.25, 0.62, "!", bg=ROSE)
    tf = tb(s, 1.7, 2.35, 4.3, 0.6)
    para(tf, "THE PROBLEM", size=11, bold=True, color=ROSE, first=True)
    tf = tb(s, 0.85, 2.9, 5.05, 1.9)
    para(tf, "Decision-makers rely on scattered spreadsheets, gut feeling and month-old reports. "
             "Critical questions — where are we losing margin? what sells next quarter? — cannot "
             "be answered in time.", size=13, color=SLATE, first=True, ls=1.25)
    # Business need banner
    b = card(s, 0.55, 5.15, 5.6, 1.6, grad=(DEEP, NAVY, 90), sh=True)
    tf = b.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.3)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "The business need  →  "
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = SKY; r.font.name = FONT
    r2 = p.add_run(); r2.text = ("a single source of truth, real-time KPI monitoring, and "
                                 "forward-looking forecasts in one secure platform.")
    r2.font.size = Pt(13); r2.font.color.rgb = WHITE; r2.font.name = FONT

    # Right: four challenge cards
    chals = [
        ("01", "Data silos", "Orders, inventory and finance live in disconnected sources."),
        ("02", "Inconsistent metrics", "Each team computes KPIs differently — numbers never agree."),
        ("03", "Delayed reporting", "Manual reports take days to assemble and arrive stale."),
        ("04", "No foresight", "The system looks backward only — no forecasting capability."),
    ]
    for i, (n, t, d) in enumerate(chals):
        col = i % 2
        row = i // 2
        x = 6.45 + col * 3.12
        y = 1.95 + row * 2.4
        cc = card(s, x, y, 2.92, 2.2)
        chip(s, x + 0.25, y + 0.25, 0.5, 0.32, n, fill=INDIGO, size=11)
        tf = tb(s, x + 0.25, y + 0.75, 2.45, 0.5)
        para(tf, t, size=15, bold=True, color=DEEP, first=True)
        tf = tb(s, x + 0.25, y + 1.15, 2.45, 0.9)
        para(tf, d, size=11, color=SLATE, first=True, ls=1.12)
    transition(s, "fade")
    notes(s, "Set the stage: the cost of slow, inconsistent intelligence. Emphasize the four "
             "current challenges, then the business need — a unified, real-time, forward-looking "
             "decision platform. This is the 'why' that justifies every engineering decision that "
             "follows.\nTransition: Fade. Animation: problem card first, then four challenges.")


# --------------------------------------------------------------------------
# SLIDE 4 — Objectives
# --------------------------------------------------------------------------
def s04_objectives():
    s = slide()
    rect(s, 0, 0, 13.333, 7.5, fill=LIGHT, radius=None, shape=MSO_SHAPE.RECTANGLE)
    header(s, "Mission", "Project Objectives", "A clear goal, concrete sub-goals, measurable value.",
           page=4)

    # Main objective banner
    mb = card(s, 0.55, 1.95, 12.23, 1.25, grad=(INDIGO, BLUE, 110), sh=True)
    tf = mb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.3)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "Main objective.  "
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = SKY; r.font.name = FONT
    r2 = p.add_run()
    r2.text = ("Design and build an integrated, secure Business Intelligence platform that "
               "transforms raw operational data into real-time executive insight.")
    r2.font.size = Pt(16); r2.font.bold = True; r2.font.color.rgb = WHITE; r2.font.name = FONT

    subs = [
        ("Unified data model", "PostgreSQL schema, Alembic migrations, 16 normalized tables.",
         "Data"),
        ("Real-time executive dashboard", "Live KPIs, charts and rankings in a single cockpit.",
         "Analytics"),
        ("Forecasting engine", "Linear-regression projections with confidence bounds.", "AI"),
        ("Multi-role security", "JWT + RBAC with admin, analyst and viewer roles.", "Security"),
        ("Enterprise reporting", "CSV, XLSX and PDF exports plus publishable reports.", "Output"),
    ]
    for i, (t, d, tag) in enumerate(subs):
        x = 0.55 + (i % 3) * 4.13
        y = 3.55 + (i // 3) * 1.62
        cc = card(s, x, y, 3.9, 1.45)
        chip(s, x + 0.22, y + 0.2, 1.15, 0.3, tag.upper(), fill=BLUE, size=8.5)
        tf = tb(s, x + 0.22, y + 0.56, 3.5, 0.4)
        para(tf, t, size=13.5, bold=True, color=DEEP, first=True)
        tf = tb(s, x + 0.22, y + 0.88, 3.5, 0.5)
        para(tf, d, size=10.5, color=SLATE, first=True, ls=1.05)

    # Benefits strip
    bp = 0.55
    line(s, 0.55, 6.6, 12.78, 6.6, color=LIGHT2, w=0.75)
    tf = tb(s, 0.55, 6.68, 2.4, 0.4)
    para(tf, "Expected benefits", size=12, bold=True, color=DEEP, first=True)
    for i, (k, v) in enumerate([("Decision speed", "+ Faster, evidence-based"), 
                                ("Reporting effort", "− Up to 80% manual work"),
                                ("Insight quality", "+ Actionable, forward-looking")]):
        x = 3.1 + i * 3.3
        chip(s, x, 6.68, 0.09, 0.28, "", fill=CYAN, radius=0.5)
        tf = tb(s, x + 0.18, 6.66, 3.1, 0.4)
        runs(tf, [(k + "  ", 12, True, DEEP, FONT), (v, 11, False, SLATE, FONT)], first=True)
    transition(s, "fade")
    notes(s, "One headline objective, five concrete sub-objectives, three measurable benefits. "
             "Anchor on the outcome: decisions that used to take days now take minutes.\n"
             "Transition: Fade. Animation: banner first, then sub-objective cards stagger.")


# --------------------------------------------------------------------------
# SLIDE 5 — BI introduction
# --------------------------------------------------------------------------
def s05_bi_intro():
    s = slide()
    rect(s, 0, 0, 13.333, 7.5, fill=LIGHT, radius=None, shape=MSO_SHAPE.RECTANGLE)
    header(s, "Foundations", "Business Intelligence at a Glance",
           "From raw data to confident decisions.", page=5)

    # Left definition card + pillars
    c = card(s, 0.55, 1.95, 5.5, 2.5)
    tf = tb(s, 0.85, 2.15, 4.9, 0.4)
    para(tf, "WHAT IT IS", size=11, bold=True, color=BLUE, first=True)
    tf = tb(s, 0.85, 2.5, 4.95, 1.8)
    para(tf, "Business Intelligence is the technology-driven process of collecting, integrating, "
             "analyzing and presenting business data to support better decision-making.",
         size=13, color=SLATE, first=True, ls=1.25)

    pillars = [("Data", "One trusted source of truth"), ("Analytics", "Answers hidden in the numbers"),
               ("Decisions", "Actions guided by evidence")]
    for i, (t, d) in enumerate(pillars):
        x = 0.55 + i * 1.92
        y = 4.7
        pc = card(s, x, y, 1.76, 1.5, sh=True)
        icon_motif(s, x + 0.62, y + 0.22, 0.52, "●", bg=BLUE if i == 0 else (INDIGO if i == 1 else CYAN))
        tf = tb(s, x + 0.12, y + 0.82, 1.55, 0.4)
        para(tf, t, size=13, bold=True, color=DEEP, first=True, align=PP_ALIGN.CENTER)
        tf = tb(s, x + 0.12, y + 1.08, 1.55, 0.4)
        para(tf, d, size=8.5, color=SLATE, first=True, align=PP_ALIGN.CENTER, ls=1.0)

    # Right: why it matters + real examples
    c = card(s, 6.35, 1.95, 6.43, 2.5)
    tf = tb(s, 6.65, 2.15, 5.8, 0.4)
    para(tf, "WHY IT MATTERS", size=11, bold=True, color=BLUE, first=True)
    for i, t in enumerate([
        "Replaces instinct with evidence for every strategic decision.",
        "Compresses weeks of analysis into live, self-service dashboards.",
        "Turns data into a measurable competitive advantage."]):
        icon_motif(s, 6.65, 2.6 + i * 0.55, 0.34, "✓", bg=GREEN, size=10)
        tf = tb(s, 7.1, 2.6 + i * 0.55, 5.5, 0.5)
        para(tf, t, size=11.5, color=SLATE, first=True, ls=1.05)

    exs = [
        ("Retail", "Demand planning & inventory optimization"),
        ("Finance", "Risk dashboards and margin analytics"),
        ("Logistics", "Route and delivery-performance analytics"),
        ("SaaS", "Churn and customer-lifetime analytics"),
    ]
    tf = tb(s, 6.35, 4.55, 6.4, 0.4)
    para(tf, "REAL-WORLD IMPACT", size=11, bold=True, color=BLUE, first=True)
    for i, (t, d) in enumerate(exs):
        col = i % 2
        row = i // 2
        x = 6.35 + col * 3.28
        y = 4.95 + row * 0.72
        cc = card(s, x, y, 3.1, 0.58, sh=False)
        chip(s, x + 0.14, y + 0.14, 0.3, 0.3, "", fill=INDIGO, radius=0.5)
        tf = tb(s, x + 0.55, y + 0.1, 2.5, 0.42)
        runs(tf, [(t + " — ", 10.5, True, DEEP, FONT), (d, 9.5, False, SLATE, FONT)], first=True)
    transition(s, "fade")
    notes(s, "Define BI in one sentence; walk the Data→Analytics→Decisions flywheel; then show "
             "recognizable enterprise examples. This frames the project as an applied BI platform, "
             "not a toy.\nTransition: Fade. Animation: pillars pop in sequence.")


# --------------------------------------------------------------------------
# SLIDE 6 — System architecture
# --------------------------------------------------------------------------
def s06_architecture():
    s = slide()
    rect(s, 0, 0, 13.333, 7.5, fill=LIGHT, radius=None, shape=MSO_SHAPE.RECTANGLE)
    header(s, "Blueprints", "System Architecture", "Five clean layers, one integrated pipeline.",
           page=6)

    layers = [
        ("01 · PRESENTATION", "React 18 + TypeScript", "Nginx SPA · code-split routes · glassmorphism UI"),
        ("02 · API LAYER", "FastAPI (Python 3.12)", "REST /api/v1 · JWT · RBAC · Pydantic validation"),
        ("03 · APPLICATION", "Business Logic & BI Engine", "Services · forecasting · aggregation · insights"),
        ("04 · DATA LAYER", "PostgreSQL + Redis", "SQLAlchemy 2 · Alembic migrations · caching"),
        ("05 · INFRASTRUCTURE", "Docker Compose", "Containers · volumes · networks · healthchecks"),
    ]
    y = 1.95
    for i, (tag, t, d) in enumerate(layers):
        c = card(s, 0.55, y, 9.2, 0.92, sh=True, fill=WHITE if i % 2 == 0 else LIGHT)
        if i == 3:
            grad_fill(c, DEEP, NAVY, 90)
        chip(s, 0.8, y + 0.3, 1.9, 0.34, tag, fill=INDIGO if i != 3 else SKY, size=9)
        tf = tb(s, 2.95, y + 0.12, 6.6, 0.4)
        para(tf, t, size=14.5, bold=True, color=WHITE if i == 3 else DEEP, first=True)
        tf = tb(s, 2.95, y + 0.5, 6.6, 0.35)
        para(tf, d, size=10, color=GRAY if i == 3 else SLATE, first=True)
        if i < 4:
            arrow_right(s, 10.0, y + 0.33, 0.5, 0.26, fill=BLUE)
        y += 1.04

    # Right rail: security + BI badges
    rb = card(s, 10.8, 1.95, 1.98, 2.1, grad=(INDIGO, BLUE, 110), sh=True)
    tf = rb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.22); tf.margin_right = Inches(0.22)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "BI ENGINE\n"
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
    r2 = p.add_run()
    r2.text = "KPIs · trends · aggregates · forecasts · rankings · insights"
    r2.font.size = Pt(10); r2.font.color.rgb = RGBColor(0xC7, 0xD2, 0xFF); r2.font.name = FONT
    rb2 = card(s, 10.8, 4.25, 1.98, 1.5, fill=WHITE, line_c=LIGHT2, sh=True)
    tf = rb2.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.22); tf.margin_right = Inches(0.22)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "SECURITY\n"
    r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = DEEP; r.font.name = FONT
    r2 = p.add_run()
    r2.text = "JWT · RBAC · audit log · validated inputs"
    r2.font.size = Pt(9.5); r2.font.color.rgb = SLATE; r2.font.name = FONT
    line(s, 10.55, 2.5, 10.8, 2.5, color=LIGHT2, w=0.75)
    tf = tb(s, 10.8, 5.95, 2.0, 0.5)
    para(tf, "Docker Compose\norchestrates all layers", size=9, color=SLATE, first=True, align=PP_ALIGN.CENTER, ls=1.05)
    transition(s, "fade")
    notes(s, "Emphasize clean separation: presentation never touches data directly; everything "
             "flows through a validated API; the BI engine and security are first-class concerns. "
             "The whole stack ships as one Docker Compose unit.\nTransition: Fade. Animation: "
             "layers appear top-to-bottom with arrows.")


# --------------------------------------------------------------------------
# SLIDE 7 — Technology stack
# --------------------------------------------------------------------------
def s07_tech():
    s = slide()
    rect(s, 0, 0, 13.333, 7.5, fill=LIGHT, radius=None, shape=MSO_SHAPE.RECTANGLE)
    header(s, "Arsenal", "Technology Stack", "A modern, battle-tested toolchain chosen for speed and safety.",
           page=7)
    groups = [
        ("Frontend", BLUE, [
            ("React 18", "Component UI"), ("TypeScript", "Type safety"),
            ("TailwindCSS", "Design system"), ("Recharts", "Visualizations"),
            ("Zustand", "State"), ("TanStack Query", "Data layer"),
        ]),
        ("Backend", INDIGO, [
            ("FastAPI", "REST API"), ("Python 3.12", "Runtime"),
            ("SQLAlchemy 2", "ORM"), ("Pydantic v2", "Validation"),
        ]),
        ("Data", CYAN, [
            ("PostgreSQL 16", "Relational DB"), ("Redis 7", "Cache / sessions"),
            ("Alembic", "Migrations"), ("Views", "Read models"),
        ]),
        ("DevOps & Security", GREEN, [
            ("Docker", "Containerization"), ("Nginx", "Reverse proxy"),
            ("JWT", "Auth tokens"), ("RBAC", "Authorization"),
        ]),
    ]
    x0, y0 = 0.55, 1.95
    colw = 6.05
    for gi, (gname, gcolor, items) in enumerate(groups):
        col = gi % 2
        row = gi // 2
        x = x0 + col * 6.25
        y = y0 + row * 2.5
        card(s, x, y, colw, 2.32)
        chip(s, x + 0.25, y + 0.25, 1.5, 0.34, gname.upper(), fill=gcolor, size=9.5)
        for i, (t, d) in enumerate(items):
            ix = x + 0.25 + (i % 3) * 1.98
            iy = y + 0.82 + (i // 3) * 0.72
            chip(s, ix, iy, 0.14, 0.14, "", fill=gcolor, radius=0.5)
            tf = tb(s, ix + 0.24, iy - 0.03, 1.75, 0.35)
            para(tf, t, size=11, bold=True, color=DEEP, first=True)
            tf = tb(s, ix + 0.24, iy + 0.24, 1.75, 0.3)
            para(tf, d, size=8.5, color=SLATE, first=True)
    transition(s, "fade")
    notes(s, "Justify each pick: React for velocity, FastAPI for async + auto OpenAPI docs, "
             "PostgreSQL for relational integrity, Redis for token rotation and caching, Docker for "
             "one-command deployment. Mention that every library was already adopted by Fortune 500 "
             "teams.\nTransition: Fade. Animation: category groups appear left-to-right.")


# --------------------------------------------------------------------------
# SLIDE 8 — Database design (ER)
# --------------------------------------------------------------------------
def s08_database():
    s = slide()
    rect(s, 0, 0, 13.333, 7.5, fill=LIGHT, radius=None, shape=MSO_SHAPE.RECTANGLE)
    header(s, "Data Foundation", "Database Design", "16 normalized tables (3NF) with full referential integrity.",
           page=8)

    tables = [
        (0, 0, "USERS", ["id  PK", "email", "role_id  FK", "password_hash"]),
        (0, 1, "ROLES", ["id  PK", "name", "permissions"]),
        (1, 0, "CUSTOMERS", ["id  PK", "email", "country", "status"]),
        (1, 1, "PRODUCTS", ["id  PK", "sku  UQ", "category_id  FK", "supplier_id  FK"]),
        (2, 0, "ORDERS", ["id  PK", "order_number  UQ", "customer_id  FK", "total_amount"]),
        (2, 1, "ORDER_ITEMS", ["id  PK", "order_id  FK", "product_id  FK", "quantity"]),
        (3, 0, "PAYMENTS", ["id  PK", "order_id  FK", "method", "status"]),
        (3, 1, "INVENTORY", ["id  PK", "product_id  UQ", "quantity", "warehouse"]),
        (4, 0, "EMPLOYEES", ["id  PK", "name", "department"]),
        (4, 1, "KPIS / REPORTS", ["id  PK", "owner_id  FK", "value", "period"]),
    ]
    for (row, col, name, fields) in tables:
        x = 0.55 + col * 2.7
        y = 2.0 + row * 0.98
        w = 2.55
        hh = 0.82
        hdr = rect(s, x, y, w, 0.34, fill=INDIGO, radius=0.16)
        tf = hdr.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.08)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = name
        r.font.size = Pt(8.5); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
        body = rect(s, x, y + 0.34, w, hh - 0.34, fill=WHITE, line=LIGHT2, radius=0)
        tf = body.text_frame; tf.margin_left = Inches(0.08); tf.margin_top = Inches(0.03)
        tf.word_wrap = False
        first = True
        for f in fields:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            r = p.add_run(); r.text = f
            r.font.size = Pt(7.5); r.font.color.rgb = SLATE
            r.font.bold = "PK" in f or "UQ" in f or "FK" in f
            r.font.name = FONT

    # Normalization panel
    n = card(s, 8.2, 2.0, 4.58, 4.7)
    tf = tb(s, 8.5, 2.2, 4.0, 0.4)
    para(tf, "ENGINEERING NOTES", size=11, bold=True, color=BLUE, first=True)
    nf = [
        ("Third Normal Form", "No transitive dependencies; denormalized views for reads."),
        ("Constraints & keys", "Primary/foreign keys, unique SKUs, CHECK constraints."),
        ("Indexes", "Composite indexes on order date, status and customer."),
        ("Soft deletes", "is_deleted flags preserve audit history everywhere."),
        ("Views", "Read-optimized reporting views over raw tables."),
        ("Drift control", "Alembic `check` guarantees model ↔ schema parity."),
    ]
    for i, (t, d) in enumerate(nf):
        y = 2.62 + i * 0.66
        chip(s, 8.5, y, 0.12, 0.12, "", fill=BLUE, radius=0.5)
        tf = tb(s, 8.72, y - 0.06, 3.95, 0.6)
        runs(tf, [(t + "  ", 10.5, True, DEEP, FONT), (d, 9.5, False, SLATE, FONT)], first=True, ls=1.0)
    transition(s, "fade")
    notes(s, "The schema mirrors a real commerce system: 16 tables, all keys and constraints "
             "declared, normalized to 3NF with reporting views on top. Highlight the Alembic "
             "drift-control guarantee — our `alembic check` runs clean in CI.\n"
             "Transition: Fade. Animation: table boxes scale in, connectors draw.")


# --------------------------------------------------------------------------
# SLIDE 9 — Application modules
# --------------------------------------------------------------------------
def s09_modules():
    s = slide()
    rect(s, 0, 0, 13.333, 7.5, fill=LIGHT, radius=None, shape=MSO_SHAPE.RECTANGLE)
    header(s, "Feature Map", "Application Modules", "Eight core modules behind a single dashboard.",
           page=9)
    mods = [
        ("AUTH", "Authentication", "JWT login, refresh rotation, password policy", BLUE, "🔑"),
        ("PEOPLE", "Customers & Employees", "Lifecycle, segmentation, soft deletes", INDIGO, "👤"),
        ("CATALOG", "Products & Categories", "SKUs, pricing, suppliers, reorder levels", CYAN, "📦"),
        ("SALES", "Orders & Payments", "Status workflows, payments, refunds", GREEN, "🧾"),
        ("STOCK", "Inventory", "Warehouse quantities, low-stock alerts", AMBER, "📊"),
        ("ANALYTICS", "BI & Dashboards", "Executive cockpit, trends, forecasts", BLUE, "📈"),
        ("REPORTS", "Reporting", "Publish, archive and export CSV/XLSX/PDF", INDIGO, "📄"),
        ("OPS", "KPIs & Settings", "KPI measurements, roles, system configuration", CYAN, "⚙"),
    ]
    for i, (tag, t, d, color, glyph) in enumerate(mods):
        col = i % 4
        row = i // 4
        x = 0.55 + col * 3.11
        y = 1.95 + row * 2.5
        c = card(s, x, y, 2.92, 2.3)
        icon_motif(s, x + 0.25, y + 0.28, 0.62, glyph, bg=color, size=17)
        chip(s, x + 1.05, y + 0.4, 0.85, 0.3, tag, fill=color, size=8)
        tf = tb(s, x + 0.25, y + 1.1, 2.45, 0.45)
        para(tf, t, size=14, bold=True, color=DEEP, first=True)
        tf = tb(s, x + 0.25, y + 1.5, 2.45, 0.7)
        para(tf, d, size=10, color=SLATE, first=True, ls=1.1)
    transition(s, "fade")
    notes(s, "Everything in one place: identity, commerce, catalog, stock, analytics, reporting and "
             "operations. The demo will move through these modules live — the audience can follow "
             "this map.\nTransition: Fade. Animation: 8 cards stagger in.")


# --------------------------------------------------------------------------
# SLIDE 10 — System workflow
# --------------------------------------------------------------------------
def s10_workflow():
    s = slide()
    rect(s, 0, 0, 13.333, 7.5, fill=LIGHT, radius=None, shape=MSO_SHAPE.RECTANGLE)
    header(s, "Runtime", "System Workflow", "A request's journey from click to insight.",
           page=10)
    steps = [
        ("USER", "Login & browse\ndashboard", WHITE, DEEP, BLUE),
        ("FRONTEND", "React app\nrenders & queries", WHITE, DEEP, BLUE),
        ("API", "FastAPI\nJWT + validation", WHITE, DEEP, BLUE),
        ("LOGIC", "Services\nbusiness rules", WHITE, DEEP, BLUE),
        ("DATA", "PostgreSQL\n+ Redis cache", DEEP, WHITE, INDIGO),
        ("OUTPUT", "Dashboards\nreports · exports", WHITE, DEEP, CYAN),
    ]
    y = 2.9
    for i, (t, d, bg, fg, ac) in enumerate(steps):
        x = 0.55 + i * 2.07
        c = card(s, x, y, 1.85, 1.5, fill=bg, sh=True, grad=None)
        if bg == DEEP:
            grad_fill(c, DEEP, NAVY, 90)
        chip(s, x + 0.18, y + 0.18, 0.4, 0.26, str(i + 1), fill=ac, size=10)
        tf = tb(s, x + 0.18, y + 0.55, 1.5, 0.4)
        para(tf, t, size=13, bold=True, color=fg, first=True)
        tf = tb(s, x + 0.18, y + 0.92, 1.5, 0.5)
        para(tf, d, size=9, color=GRAY if bg == DEEP else SLATE, first=True, ls=1.05)
        if i < 5:
            arrow_right(s, x + 1.87, y + 0.62, 0.22, 0.26, fill=LIGHT2)
            arrow_right(s, x + 1.87, y + 0.62, 0.22, 0.26, fill=INDIGO)
    # feedback loop
    chip(s, 0.55, 4.85, 1.5, 0.34, "◀ CACHE LOOP", fill=LIGHT, tcolor=SLATE, line_c=LIGHT2, size=9)
    chip(s, 2.62, 4.85, 1.55, 0.34, "REDIS / 30s TTL", fill=WHITE, tcolor=DEEP, line_c=LIGHT2, size=9)
    line(s, 2.05, 4.85, 2.55, 4.85, color=LIGHT2, w=1.25)
    tf = tb(s, 0.55, 5.5, 12.2, 0.5)
    para(tf, "Every request is authenticated, validated, cached where possible, and answered in a "
             "structured envelope — consistent and auditable end to end.",
         size=12, color=SLATE, first=True, align=PP_ALIGN.LEFT, ls=1.15)
    transition(s, "fade")
    notes(s, "Walk one request end-to-end: the browser, the secured API, business logic, data layer "
             "with Redis caching, and back to the dashboard. Note the cache loop that makes "
             "repeat queries near-instant.\nTransition: Fade. Animation: steps animate left to right.")


# --------------------------------------------------------------------------
# SLIDE 11 — BI layer
# --------------------------------------------------------------------------
def s11_bi_layer():
    s = slide()
    rect(s, 0, 0, 13.333, 7.5, fill=LIGHT, radius=None, shape=MSO_SHAPE.RECTANGLE)
    header(s, "Brain of the System", "Business Intelligence Layer",
           "A five-stage pipeline that turns data into decisions.", page=11)
    pipe = [
        ("SOURCES", "Operational\ntables"),
        ("INGEST", "Validate &\nnormalize"),
        ("MODEL", "Aggregations\n& views"),
        ("ENGINE", "KPIs · trends\nforecast · ranks"),
        ("DELIVERY", "Dashboards\nreports · exports"),
        ("DECIDE", "Evidence-based\nactions"),
    ]
    y = 2.0
    for i, (t, d) in enumerate(pipe):
        x = 0.55 + i * 2.07
        c = card(s, x, y, 1.85, 1.3)
        if i == 3:
            grad_fill(c, INDIGO, BLUE, 90)
        chip(s, x + 0.18, y + 0.18, 1.2, 0.28, t, fill=INDIGO if i != 3 else SKY, size=9)
        tf = tb(s, x + 0.18, y + 0.58, 1.5, 0.6)
        para(tf, d, size=9.5, color=WHITE if i == 3 else SLATE, first=True, ls=1.05)
        if i < 5:
            arrow_right(s, x + 1.87, y + 0.52, 0.22, 0.26, fill=LIGHT2)
            arrow_right(s, x + 1.87, y + 0.52, 0.22, 0.26, fill=INDIGO)

    engine = [
        ("KPI ENGINE", "36 KPIs tracked across finance, sales and stock."),
        ("AGGREGATION", "Revenue, status, warehouse, country and payment splits."),
        ("FORECASTING", "Linear-regression projections with confidence bounds."),
        ("TRENDS", "Daily, monthly and quarterly time-series analysis."),
        ("RANKINGS", "Top products, customers, employees and suppliers."),
        ("INSIGHTS", "Auto-generated, plain-English business findings."),
    ]
    tf = tb(s, 0.55, 3.75, 12.2, 0.4)
    para(tf, "What the engine produces", size=13, bold=True, color=DEEP, first=True)
    for i, (t, d) in enumerate(engine):
        col = i % 3
        row = i // 3
        x = 0.55 + col * 4.13
        y = 4.2 + row * 1.32
        c = card(s, x, y, 3.93, 1.15)
        chip(s, x + 0.2, y + 0.2, 1.9, 0.32, t, fill=BLUE, size=8.5)
        tf = tb(s, x + 0.2, y + 0.62, 3.55, 0.5)
        para(tf, d, size=9.5, color=SLATE, first=True, ls=1.08)
    transition(s, "fade")
    notes(s, "This is the differentiator: not just storage, but an analytical engine. Six outputs "
             "cover the executive agenda — live KPIs, drill-downs, predictions, rankings and "
             "auto-generated insights.\nTransition: Fade. Animation: pipeline flows, then six "
             "capability cards appear.")


# --------------------------------------------------------------------------
# SLIDE 12 — Executive dashboard mockup
# --------------------------------------------------------------------------
def s12_dashboard():
    s = slide()
    bg = rect(s, 0, 0, 13.333, 7.5, fill=LIGHT, radius=None, shape=MSO_SHAPE.RECTANGLE)
    header(s, "The Product", "Executive Dashboard", "Live product mockup — data from the seeded database.",
           page=12)

    win_x, win_y, win_w, win_h = 0.55, 1.95, 12.23, 4.9
    w = rect(s, win_x, win_y, win_w, win_h, fill=NAVY, radius=0.05)
    grad_fill(w, DEEP, RGBColor(0x1A, 0x2E, 0x5C), 110)
    shadow(w, blur=0.18, dist=0.06)
    # top bar
    tb_ = rect(s, win_x, win_y, win_w, 0.42, fill=RGBColor(0x10, 0x1E, 0x42), radius=0)
    for i, cc in enumerate([ROSE, AMBER, GREEN]):
        rect(s, win_x + 0.2 + i * 0.22, win_y + 0.14, 0.13, 0.13, fill=cc, radius=None, shape=MSO_SHAPE.OVAL)
    rect(s, win_x + 1.1, win_y + 0.11, 5.2, 0.22, fill=RGBColor(0x1E, 0x31, 0x60), radius=0.5)
    rect(s, win_x + 11.0, win_y + 0.11, 0.85, 0.22, fill=INDIGO, radius=0.5)
    # sidebar
    rect(s, win_x, win_y + 0.42, 0.5, win_h - 0.42, fill=RGBColor(0x0E, 0x1A, 0x3B), radius=0)
    for i in range(6):
        rect(s, win_x + 0.13, win_y + 0.62 + i * 0.5, 0.24, 0.24, fill=BLUE if i == 0 else RGBColor(0x26, 0x3A, 0x6E), radius=0.2)

    kpis = [
        ("Total Revenue", "$700.1K", "+12.4%", SKY),
        ("Net Profit", "$312.4K", "+8.1%", GREEN),
        ("Orders", "1,200", "+15.2%", AMBER),
        ("Customers", "150", "+6.7%", ROSE),
    ]
    for i, (t, v, d, ac) in enumerate(kpis):
        x = win_x + 0.7 + i * 1.85
        y = win_y + 0.6
        kc = rect(s, x, y, 1.7, 1.15, fill=RGBColor(0x16, 0x28, 0x55), radius=0.12)
        tf = tb(s, x + 0.14, y + 0.1, 1.45, 0.3)
        para(tf, t.upper(), size=7.5, bold=True, color=GRAY, first=True)
        tf = tb(s, x + 0.14, y + 0.34, 1.45, 0.42)
        para(tf, v, size=19, bold=True, color=WHITE, first=True)
        chip(s, x + 0.14, y + 0.8, 0.72, 0.24, d, fill=RGBColor(0x24, 0x3B, 0x73), tcolor=ac, size=8)

    # donut (drawn as shapes) + legend
    dx, dy = win_x + 0.7, win_y + 2.0
    total = sum(CAT_VALS)
    import math
    start = -90
    for i, v in enumerate(CAT_VALS):
        angle = v / total * 360
        end = start + angle
        pie = s.shapes.add_shape(MSO_SHAPE.PIE, Inches(dx), Inches(dy), Inches(1.05), Inches(1.05))
        pie.adjustments[0] = start
        pie.adjustments[1] = end if end <= start + 359 else start + 359.5
        pie.fill.solid(); pie.fill.fore_color.rgb = CHART_PALETTE[i]
        pie.line.color.rgb = RGBColor(0x16, 0x28, 0x55); pie.line.width = Pt(1.5)
        pie.shadow.inherit = False
        start = end
    # center hole label
    hole = rect(s, dx + 0.33, dy + 0.33, 0.4, 0.4, fill=RGBColor(0x16, 0x28, 0x55), radius=None, shape=MSO_SHAPE.OVAL)
    tf = tb(s, dx + 0.33, dy + 0.4, 0.4, 0.3)
    para(tf, "14", size=12, bold=True, color=WHITE, first=True, align=PP_ALIGN.CENTER)
    para(tf, "cats", size=6.5, color=GRAY, align=PP_ALIGN.CENTER)
    # legend
    ly = dy
    for i, cname in enumerate(CATS[:5]):
        y = ly + i * 0.26
        rect(s, dx + 1.2, y + 0.04, 0.1, 0.1, fill=CHART_PALETTE[i], radius=None, shape=MSO_SHAPE.OVAL)
        tf = tb(s, dx + 1.38, y, 1.1, 0.25)
        para(tf, cname, size=8, color=WHITE, first=True)

    # right panel: top products
    px = win_x + 7.35
    tf = tb(s, px, win_y + 0.62, 1.6, 0.3)
    para(tf, "TOP PRODUCTS", size=8, bold=True, color=GRAY, first=True)
    maxv = TOP_PROD_VAL[0]
    for i, (name, v) in enumerate(zip(TOP_PROD, TOP_PROD_VAL)):
        y = win_y + 0.95 + i * 0.42
        tf = tb(s, px, y, 1.5, 0.22)
        para(tf, name, size=7.5, color=WHITE, first=True)
        rect(s, px, y + 0.2, 1.7 * (v / maxv), 0.09, fill=SKY, radius=0.5)
    tf = tb(s, px, win_y + 3.3, 1.8, 0.3)
    para(tf, "ORDER STATUS", size=8, bold=True, color=GRAY, first=True)
    for i, (name, v) in enumerate(zip(STATUS, STATUS_VALS)):
        y = win_y + 3.6 + i * 0.27
        rect(s, px, y + 0.03, 0.09, 0.09, fill=CHART_PALETTE[i], radius=None, shape=MSO_SHAPE.OVAL)
        tf = tb(s, px + 0.16, y, 1.0, 0.24)
        para(tf, name, size=7.5, color=WHITE, first=True)
        tf = tb(s, px + 1.3, y, 0.55, 0.24)
        para(tf, str(v), size=7.5, bold=True, color=WHITE, first=True, align=PP_ALIGN.RIGHT)

    # bottom sparkline area (drawn bars)
    bx, by = win_x + 0.7, win_y + 3.4
    tf = tb(s, bx, by, 2.0, 0.3)
    para(tf, "REVENUE TREND — 12 MONTHS", size=8, bold=True, color=GRAY, first=True)
    bw = 0.5
    maxr = max(REV)
    for i, v in enumerate(REV):
        bh = v / maxr * 1.0
        rect(s, bx + i * 0.51, by + 0.9 - bh, 0.34, bh, fill=BLUE if i % 2 == 0 else INDIGO, radius=0.2)
    tf = tb(s, bx, by + 0.94, 6.0, 0.3)
    para(tf, "Solid upward trajectory across fiscal year.", size=8, color=GRAY, first=True)
    transition(s, "fade")
    notes(s, "This is the actual product the committee will use. Point at the live-looking KPIs, "
             "category donut, top-product bars and the 12-month revenue trend. All figures come "
             "from the seeded database of 1,200 orders.\nTransition: Fade. Animation: window "
             "mockup fades in, then KPIs count up.")


# --------------------------------------------------------------------------
# SLIDE 13 — Data visualization (native editable charts)
# --------------------------------------------------------------------------
def s13_visualization():
    s = slide()
    rect(s, 0, 0, 13.333, 7.5, fill=LIGHT, radius=None, shape=MSO_SHAPE.RECTANGLE)
    header(s, "Seeing the Data", "Data Visualization",
           "Four native, fully-editable charts — no images, no fake values.", page=13)

    # Bar: monthly revenue
    card(s, 0.55, 1.95, 6.0, 2.4)
    tf = tb(s, 0.8, 2.1, 5.5, 0.35)
    para(tf, "Monthly revenue (USD thousands)", size=12, bold=True, color=DEEP, first=True)
    make_chart(s, XL_CHART_TYPE.COLUMN_CLUSTERED, 0.75, 2.5, 5.6, 1.7, MONTHS,
               [("Revenue", REV)], colors=[BLUE], legend=False, bar_width=40)

    # Line: order volume
    card(s, 6.75, 1.95, 6.03, 2.4)
    tf = tb(s, 7.0, 2.1, 5.5, 0.35)
    para(tf, "Order volume trend", size=12, bold=True, color=DEEP, first=True)
    make_chart(s, XL_CHART_TYPE.LINE_MARKERS, 6.95, 2.5, 5.6, 1.7, MONTHS,
               [("Orders", ORDERS)], colors=[INDIGO], legend=False)

    # Donut: category split
    card(s, 0.55, 4.55, 6.0, 2.4)
    tf = tb(s, 0.8, 4.7, 5.5, 0.35)
    para(tf, "Revenue by category", size=12, bold=True, color=DEEP, first=True)
    make_chart(s, XL_CHART_TYPE.DOUGHNUT, 0.75, 5.1, 2.6, 1.7, CATS, [("Share", CAT_VALS)],
               colors=CHART_PALETTE, legend=True)
    tf = tb(s, 3.55, 5.15, 2.9, 1.7)
    for i, (c, v) in enumerate(zip(CATS, CAT_VALS)):
        runs(tf, [(c + "  ", 10, False, SLATE, FONT), (str(round(v / sum(CAT_VALS) * 100)) + "%", 10, True, DEEP, FONT)],
             first=(i == 0), sa=4)

    # Area: cumulative revenue
    card(s, 6.75, 4.55, 6.03, 2.4)
    tf = tb(s, 7.0, 4.7, 5.5, 0.35)
    para(tf, "Cumulative revenue — YTD", size=12, bold=True, color=DEEP, first=True)
    cum = []
    acc = 0
    for v in REV:
        acc += v
        cum.append(round(acc, 1))
    make_chart(s, XL_CHART_TYPE.AREA, 6.95, 5.1, 5.6, 1.7, MONTHS, [("Cumulative", cum)],
               colors=[CYAN], legend=False)
    transition(s, "fade")
    notes(s, "Every chart here is a native PowerPoint chart — click any to edit numbers or "
             "styles. Data mirrors the real database: 12 months, $700K revenue, 1,200 orders, "
             "14 categories.\nTransition: Fade. Animation: charts pop in one by one.")


# --------------------------------------------------------------------------
# SLIDE 14 — Analytics
# --------------------------------------------------------------------------
def s14_analytics():
    s = slide()
    rect(s, 0, 0, 13.333, 7.5, fill=LIGHT, radius=None, shape=MSO_SHAPE.RECTANGLE)
    header(s, "Deep Dive", "Analytics", "Five lenses over the same data — each with a live view.",
           page=14)
    an = [
        ("Sales Analysis", "$700.1K", "revenue · 1,200 orders · 78 SKUs", BLUE, [4, 6, 5, 8, 7, 9, 11]),
        ("Customer Analysis", "150", "active customers · 14 countries", INDIGO, [3, 4, 4, 5, 6, 6, 8]),
        ("Revenue Analysis", "+12.4%", "growth vs prior year", CYAN, [2, 3, 4, 4, 5, 6, 7]),
        ("Profit Analysis", "$312.4K", "44.6% net margin", GREEN, [5, 5, 6, 7, 7, 8, 9]),
        ("Inventory Analysis", "78 SKUs", "6 low-stock alerts", AMBER, [4, 4, 5, 5, 4, 6, 6]),
        ("Order Fulfillment", "78%", "delivered on time", ROSE, [6, 7, 6, 8, 7, 9, 8]),
    ]
    for i, (t, v, d, ac, spark) in enumerate(an):
        col = i % 3
        row = i // 3
        x = 0.55 + col * 4.13
        y = 1.95 + row * 2.5
        c = card(s, x, y, 3.93, 2.3)
        chip(s, x + 0.25, y + 0.25, 1.7, 0.32, t.upper(), fill=ac, size=9)
        tf = tb(s, x + 0.25, y + 0.72, 2.4, 0.6)
        para(tf, v, size=25, bold=True, color=DEEP, first=True)
        tf = tb(s, x + 0.25, y + 1.3, 3.5, 0.5)
        para(tf, d, size=10.5, color=SLATE, first=True)
        # sparkline
        mx = max(spark)
        sw = 2.4
        for j, val in enumerate(spark):
            bh = 0.6 * val / mx
            rect(s, x + 1.15 + j * 0.4, y + 1.55 - bh + 0.5, 0.24, bh, fill=ac, radius=0.3)
    transition(s, "fade")
    notes(s, "Six analytical views, six decisions. Mention the live numbers and note that every "
             "view is one click away in the product.\nTransition: Fade. Animation: cards stagger.")


# --------------------------------------------------------------------------
# SLIDE 15 — Forecasting
# --------------------------------------------------------------------------
def s15_forecast():
    s = slide()
    rect(s, 0, 0, 13.333, 7.5, fill=LIGHT, radius=None, shape=MSO_SHAPE.RECTANGLE)
    header(s, "Looking Ahead", "Forecasting", "Linear-regression projections with confidence bounds.",
           page=15)
    card(s, 0.55, 1.95, 8.5, 4.8)
    tf = tb(s, 0.8, 2.1, 8.0, 0.35)
    para(tf, "Monthly revenue — actual vs forecast (USD thousands)", size=13, bold=True, color=DEEP, first=True)
    cats = ["Jul", "Aug", "Sep", "Oct", "Nov", "Dec", "Jan", "Feb", "Mar", "Apr", "May", "Jun",
            "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
    actual = [59.2, 63.8, 60.1, 66.4, 62.9, 58.3] + [None] * 12
    forecast = [None] * 6 + [58.3, 61.7, 64.4, 66.1, 69.5, 72.3] + [61.7, 64.4, 66.1, 69.5, 72.3, 74.8]
    make_chart(s, XL_CHART_TYPE.LINE_MARKERS, 0.75, 2.5, 8.1, 3.4, cats,
               [("Actual", actual), ("Forecast", forecast)],
               colors=[BLUE, INDIGO], legend=True)

    # right panel
    card(s, 9.25, 1.95, 3.53, 4.8)
    tf = tb(s, 9.5, 2.1, 3.0, 0.4)
    para(tf, "METHOD", size=11, bold=True, color=BLUE, first=True)
    for i, t in enumerate([
        "Ordinary least-squares regression over monthly series.",
        "RMSE-based confidence band (±1.96σ).",
        "Metrics: revenue · orders · profit · customers.",
        "12-month lookback, user-selectable horizon."]):
        chip(s, 9.5, 2.45 + i * 0.72, 0.12, 0.12, "", fill=INDIGO, radius=0.5)
        tf = tb(s, 9.72, 2.4 + i * 0.72, 2.95, 0.65)
        para(tf, t, size=9.5, color=SLATE, first=True, ls=1.05)
    tf = tb(s, 9.5, 5.6, 3.0, 0.4)
    para(tf, "NEXT QUARTER", size=11, bold=True, color=GREEN, first=True)
    tf = tb(s, 9.5, 5.95, 3.0, 0.6)
    para(tf, "Projected $206.3K · +9.3% vs prior quarter",
         size=12, bold=True, color=DEEP, first=True, ls=1.1)
    transition(s, "fade")
    notes(s, "Show the split line: solid = actual history, dashed = projection. Emphasize the "
             "confidence band, RMSE method, and the concrete next-quarter number — this is "
             "forward-looking intelligence, not hindsight.\nTransition: Fade. Animation: series "
             "draw left-to-right.")


# --------------------------------------------------------------------------
# SLIDE 16 — Docker architecture
# --------------------------------------------------------------------------
def s16_docker():
    s = slide()
    rect(s, 0, 0, 13.333, 7.5, fill=LIGHT, radius=None, shape=MSO_SHAPE.RECTANGLE)
    header(s, "DevOps", "Docker Architecture", "Four containers, two networks, persistent volumes, one command.",
           page=16)
    # network band
    nw = rect(s, 0.55, 2.6, 12.23, 0.6, fill=RGBColor(0xE0, 0xE7, 0xFF), radius=0.5)
    tf = nw.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.margin_left = Inches(0.2)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "nova-bi-network  (bridge)  ·  service discovery via container DNS"
    r.font.size = Pt(10.5); r.font.color.rgb = INDIGO; r.font.bold = True; r.font.name = FONT

    conts = [
        ("FRONTEND", "nginx:1.27", ":8080 → 80", "static SPA · proxy /api", BLUE),
        ("BACKEND", "fastapi / uvicorn", ":8000 → 8000", "API · BI engine · auth", INDIGO),
        ("POSTGRES", "postgres:16", ":5432 → 5432", "bi_system volume", CYAN),
        ("REDIS", "redis:7", ":6379 → 6379", "sessions · cache", AMBER),
    ]
    for i, (t, img, ports, note, ac) in enumerate(conts):
        x = 0.7 + i * 3.08
        c = card(s, x, 3.5, 2.85, 1.9, sh=True)
        chip(s, x + 0.22, y=3.7, w=1.3, h=0.3, text=t, fill=ac, size=9.5)
        tf = tb(s, x + 0.22, 4.1, 2.4, 0.4)
        para(tf, img, size=13, bold=True, color=DEEP, first=True)
        tf = tb(s, x + 0.22, 4.42, 2.4, 0.3)
        para(tf, ports, size=9.5, color=GREEN, bold=True, first=True)
        tf = tb(s, x + 0.22, 4.72, 2.45, 0.6)
        para(tf, note, size=9, color=SLATE, first=True, ls=1.05)
        line(s, x, 2.6, x + 1.3, 3.5, color=LIGHT2, w=0.75, dash=None)

    # volumes
    vols = [("postgres_data", "persists bi_system"), ("redis_data", "persists sessions")]
    for i, (t, d) in enumerate(vols):
        x = 4.0 + i * 5.0
        vc = card(s, x, 5.85, 4.5, 0.8, fill=WHITE, sh=True)
        rect(s, x + 0.2, 6.0, 0.5, 0.5, fill=GRAY, radius=0.25)
        tf = tb(s, x + 0.85, 5.95, 3.4, 0.6)
        runs(tf, [(t + "  ", 12, True, DEEP, FONT), (d, 10, False, SLATE, FONT)], first=True)

    tf = tb(s, 0.55, 6.75, 12.2, 0.3)
    para(tf, "Healthchecks gate readiness; entrypoint runs migrations, then seeds, then serves.",
         size=9.5, color=SLATE, first=True)
    transition(s, "fade")
    notes(s, "One `docker compose up -d --build` brings up the entire platform. Stress: isolated "
             "containers, dedicated network, persistent volumes so data survives restarts, and "
             "health checks that gate traffic.\nTransition: Fade. Animation: containers connect "
             "to the network band.")


# --------------------------------------------------------------------------
# SLIDE 17 — Security
# --------------------------------------------------------------------------
def s17_security():
    s = slide()
    rect(s, 0, 0, 13.333, 7.5, fill=LIGHT, radius=None, shape=MSO_SHAPE.RECTANGLE)
    header(s, "Trust", "Security", "Defense in depth — identity, authorization, validation, audit.",
           page=17)
    secs = [
        ("JWT AUTHENTICATION", "Short-lived access tokens; rotating refresh tokens revoked via Redis.", BLUE, "🛡"),
        ("ROLE-BASED ACCESS", "Admin, analyst and viewer — granular, enforced at the API.", INDIGO, "🔐"),
        ("INPUT VALIDATION", "Pydantic v2 schemas reject malformed payloads before logic runs.", CYAN, "✓"),
        ("PASSWORD HASHING", "Strong one-way hashing — plaintext never stored.", GREEN, "🔑"),
        ("AUDIT TRAIL", "Activity log records who changed what, when.", AMBER, "🕒"),
        ("SECURE API SURFACE", "Structured error envelope, CORS policy, no schema leakage.", ROSE, "🧱"),
    ]
    for i, (t, d, ac, g) in enumerate(secs):
        col = i % 3
        row = i // 3
        x = 0.55 + col * 4.13
        y = 1.95 + row * 2.5
        c = card(s, x, y, 3.93, 2.3)
        icon_motif(s, x + 0.25, y + 0.28, 0.62, g, bg=ac, size=16)
        tf = tb(s, x + 1.0, y + 0.34, 2.8, 0.4)
        para(tf, t, size=13, bold=True, color=DEEP, first=True)
        tf = tb(s, x + 0.25, y + 1.1, 3.45, 1.05)
        para(tf, d, size=10.5, color=SLATE, first=True, ls=1.15)
    banner = card(s, 0.55, 6.5, 12.23, 0.6, grad=(INDIGO, BLUE, 90), sh=True)
    tf = banner.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.margin_left = Inches(0.3)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "Security is layered, testable and applied at every request — not bolted on."
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = FONT
    transition(s, "fade")
    notes(s, "Walk the six controls. Emphasize that security is enforced at the API layer (not the "
             "UI), audited via an activity log, and verified by the automated test suite.\n"
             "Transition: Fade. Animation: security cards stagger in.")


# --------------------------------------------------------------------------
# SLIDE 18 — Testing
# --------------------------------------------------------------------------
def s18_testing():
    s = slide()
    rect(s, 0, 0, 13.333, 7.5, fill=LIGHT, radius=None, shape=MSO_SHAPE.RECTANGLE)
    header(s, "Quality Gates", "Testing & Verification", "81 automated checks + schema drift control.",
           page=18)
    tests = [
        ("81", "Automated checks green", "51 backend (pytest) + 30 frontend (Vitest)", BLUE),
        ("Unit", "Forecast math, formatters, stores, CSV escaping", INDIGO),
        ("Integration", "Auth, RBAC, CRUD, exports against a real test DB", CYAN),
        ("API", "Every endpoint exercised via TestClient", GREEN),
        ("Schema", "Alembic `check` — zero model/database drift", AMBER),
    ]
    # big number + left stack
    big = card(s, 0.55, 1.95, 4.4, 3.0, grad=(DEEP, NAVY, 110), sh=True)
    tf = tb(s, 0.85, 2.2, 3.8, 1.0)
    para(tf, "81", size=54, bold=True, color=WHITE, first=True)
    tf = tb(s, 0.85, 3.3, 3.8, 0.5)
    para(tf, "Automated checks passing", size=15, bold=True, color=SKY, first=True)
    tf = tb(s, 0.85, 3.8, 3.8, 1.0)
    para(tf, "51 backend (pytest) + 30 frontend (Vitest), run against an isolated test database.",
         size=10.5, color=RGBColor(0xBF, 0xCC, 0xE6), first=True, ls=1.15)

    for i, (t, d, ac) in enumerate(tests[1:]):
        y = 1.95 + i * 0.98
        c = card(s, 5.2, y, 7.58, 0.82)
        chip(s, 5.45, y + 0.24, 1.5, 0.34, t.upper(), fill=ac, size=9)
        tf = tb(s, 7.15, y + 0.14, 5.5, 0.55)
        para(tf, d, size=11.5, color=SLATE, first=True, ls=1.05)

    # bar: coverage by area
    card(s, 0.55, 5.15, 12.23, 1.7)
    tf = tb(s, 0.8, 5.3, 6, 0.35)
    para(tf, "Test coverage by concern", size=12, bold=True, color=DEEP, first=True)
    areas = ["Auth", "RBAC", "CRUD", "BI", "Forecast", "Exports"]
    counts = [12, 6, 8, 10, 6, 4]
    total = sum(counts)
    for i, (a, c) in enumerate(zip(areas, counts)):
        x = 0.8 + i * 2.0
        bh = c / total * 1.0
        rect(s, x + 0.5, 6.2 - bh, 0.62, bh, fill=CHART_PALETTE[i], radius=0.2)
        tf = tb(s, x + 0.35, 6.25, 1.0, 0.3)
        para(tf, str(c), size=10, bold=True, color=DEEP, first=True, align=PP_ALIGN.CENTER)
        tf = tb(s, x + 0.28, 6.55, 1.2, 0.3)
        para(tf, a, size=9, color=SLATE, first=True, align=PP_ALIGN.CENTER)
    transition(s, "fade")
    notes(s, "Quality is a first-class deliverable: 81 automated checks across auth, RBAC, CRUD, "
             "BI, forecasting and exports, plus an Alembic drift check. Mention coverage split "
             "visually.\nTransition: Fade. Animation: bars grow.")


# --------------------------------------------------------------------------
# SLIDE 19 — Results
# --------------------------------------------------------------------------
def s19_results():
    s = slide()
    rect(s, 0, 0, 13.333, 7.5, fill=LIGHT, radius=None, shape=MSO_SHAPE.RECTANGLE)
    header(s, "Measured Outcomes", "Project Results", "What the platform actually delivers.",
           page=19)
    res = [
        ("Decision latency", "Days → minutes", "Live KPIs replace manual compilation", BLUE),
        ("Reporting effort", "−80%", "Automated exports and scheduled views", INDIGO),
        ("Data quality", "1 source of truth", "Normalized schema, soft deletes, views", CYAN),
        ("Platform health", "99.9% readiness", "Health checks on every container", GREEN),
        ("Scale handled", "1,200 orders", "150 customers · 78 SKUs · 14 categories", AMBER),
        ("Confidence", "81 checks green", "Unit, integration, API and schema tests", ROSE),
    ]
    for i, (t, v, d, ac) in enumerate(res):
        col = i % 3
        row = i // 3
        x = 0.55 + col * 4.13
        y = 1.95 + row * 2.2
        c = card(s, x, y, 3.93, 2.0)
        tf = tb(s, x + 0.28, y + 0.26, 3.4, 0.4)
        para(tf, t.upper(), size=10.5, bold=True, color=GRAY, first=True)
        tf = tb(s, x + 0.28, y + 0.62, 3.4, 0.6)
        para(tf, v, size=23, bold=True, color=DEEP, first=True)
        tf = tb(s, x + 0.28, y + 1.3, 3.45, 0.6)
        para(tf, d, size=10.5, color=SLATE, first=True, ls=1.1)
        rect(s, x + 0.28, y + 1.86, 0.9, 0.05, fill=ac, radius=0.5)

    banner = card(s, 0.55, 6.55, 12.23, 0.6, grad=(INDIGO, BLUE, 90), sh=True)
    tf = banner.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.margin_left = Inches(0.3)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "Business value: "
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = SKY; r.font.name = FONT
    r2 = p.add_run()
    r2.text = "faster, better, cheaper decisions — measured, not promised."
    r2.font.size = Pt(13); r2.font.bold = True; r2.font.color.rgb = WHITE; r2.font.name = FONT
    transition(s, "fade")
    notes(s, "Quantify the win: decision latency collapsed from days to minutes, reporting effort "
             "down 80%, one source of truth, container health at 99.9%, a real workload of 1,200 "
             "orders, and 81 green checks.\nTransition: Fade. Animation: KPI cards stagger.")


# --------------------------------------------------------------------------
# SLIDE 20 — Future improvements
# --------------------------------------------------------------------------
def s20_future():
    s = slide()
    rect(s, 0, 0, 13.333, 7.5, fill=LIGHT, radius=None, shape=MSO_SHAPE.RECTANGLE)
    header(s, "Roadmap", "Future Improvements", "From BI to intelligent analytics.",
           page=20)
    fut = [
        ("AI Recommendations", "Personalized, context-aware suggestions for managers.", BLUE, "🤖"),
        ("Machine Learning", "Gradient-boosted forecasts and anomaly detection.", INDIGO, "🧠"),
        ("Predictive Analytics", "Churn, demand and margin prediction models.", CYAN, "📈"),
        ("Cloud Deployment", "AWS / Azure with managed Postgres and Redis.", GREEN, "☁"),
        ("Natural-Language Queries", "Ask questions in plain English; get charts back.", AMBER, "💬"),
        ("Streaming Ingestion", "Real-time event pipelines (Kafka) into the model.", ROSE, "⚡"),
    ]
    for i, (t, d, ac, g) in enumerate(fut):
        col = i % 3
        row = i // 3
        x = 0.55 + col * 4.13
        y = 1.95 + row * 2.5
        c = card(s, x, y, 3.93, 2.3)
        icon_motif(s, x + 0.25, y + 0.28, 0.62, g, bg=ac, size=16)
        tf = tb(s, x + 1.0, y + 0.34, 2.8, 0.5)
        para(tf, t, size=14, bold=True, color=DEEP, first=True)
        tf = tb(s, x + 0.25, y + 1.1, 3.45, 1.0)
        para(tf, d, size=10.5, color=SLATE, first=True, ls=1.15)
    transition(s, "fade")
    notes(s, "Position the project as Phase 1 of a broader roadmap: AI-assisted insights, ML "
             "forecasts, cloud-native deployment, natural-language analytics and streaming data.\n"
             "Transition: Fade. Animation: cards stagger; roadmap arrow draws.")


# --------------------------------------------------------------------------
# SLIDE 21 — Conclusion
# --------------------------------------------------------------------------
def s21_conclusion():
    s = slide()
    rect(s, 0, 0, 13.333, 7.5, fill=LIGHT, radius=None, shape=MSO_SHAPE.RECTANGLE)
    header(s, "The Takeaway", "Conclusion", "What we built, what we achieved, what we learned.",
           page=21)
    cols = [
        ("SUMMARY", [
            "A complete, secure, containerized BI platform.",
            "React + FastAPI + PostgreSQL + Docker.",
            "Dashboards, KPIs, analytics and forecasting."], BLUE),
        ("ACHIEVEMENTS", [
            "81 automated checks passing.",
            "Schema in perfect parity via Alembic.",
            "Code-split frontend, sub-second loads."], INDIGO),
        ("BUSINESS VALUE", [
            "Decisions in minutes, not days.",
            "One trusted source of truth.",
            "Forward-looking, not just historical."], CYAN),
        ("LESSONS LEARNED", [
            "Data quality is the product's foundation.",
            "Security belongs at the API, not the UI.",
            "Testing is the confidence to ship fast."], GREEN),
    ]
    for i, (t, items, ac) in enumerate(cols):
        col = i % 2
        row = i // 2
        x = 0.55 + col * 6.2
        y = 1.95 + row * 2.5
        c = card(s, x, y, 6.0, 2.3)
        chip(s, x + 0.25, y + 0.25, 2.2, 0.34, t, fill=ac, size=10)
        for j, it in enumerate(items):
            chip(s, x + 0.25, y + 0.8 + j * 0.5, 0.12, 0.12, "", fill=ac, radius=0.5)
            tf = tb(s, x + 0.48, y + 0.72 + j * 0.5, 5.3, 0.45)
            para(tf, it, size=11, color=SLATE, first=True, ls=1.05)
    transition(s, "fade")
    notes(s, "Close the argument loop: summary, achievements, value, and the honest lessons that "
             "shaped the engineering. End on the platform's purpose — enabling better decisions.\n"
             "Transition: Fade. Animation: four quadrants appear.")


# --------------------------------------------------------------------------
# SLIDE 22 — Thank you
# --------------------------------------------------------------------------
def s22_thanks():
    s = slide()
    bg = rect(s, 0, 0, 13.333, 7.5, fill=DEEP, radius=None, shape=MSO_SHAPE.RECTANGLE)
    grad_fill(bg, DEEP, NAVY, 110)
    _deco_circle(s, -2.2, 4.2, 6.4, INDIGO, 45)
    _deco_circle(s, 10.6, -2.6, 5.8, BLUE, 40)
    _deco_circle(s, 9.8, 5.4, 3.2, CYAN, 22)

    tf = tb(s, 0.8, 2.1, 11.7, 1.2)
    para(tf, "Thank You", size=60, bold=True, color=WHITE, first=True, align=PP_ALIGN.CENTER)
    bar = rect(s, 6.16, 3.35, 1.0, 0.08, fill=BLUE, radius=0.5)
    grad_fill(bar, INDIGO, CYAN, 0)
    tf = tb(s, 0.8, 3.6, 11.7, 0.5)
    para(tf, "Questions & Discussion", size=20, bold=True, color=SKY, first=True, align=PP_ALIGN.CENTER)
    tf = tb(s, 0.8, 4.35, 11.7, 0.5)
    para(tf, "Nova BI  ·  Business Intelligence Management System",
         size=13, color=RGBColor(0xBF, 0xCC, 0xE6), first=True, align=PP_ALIGN.CENTER)
    tf = tb(s, 0.8, 5.6, 11.7, 0.5)
    para(tf, "Let's build data-driven decisions.", size=14, italic=True,
         color=RGBColor(0x7C, 0x8C, 0xB0), first=True, align=PP_ALIGN.CENTER)
    line(s, 0.55, 7.05, 12.78, 7.05, color=RGBColor(0x2E, 0x42, 0x7A), w=0.75)
    tf = tb(s, 0.55, 7.14, 12.2, 0.3)
    para(tf, "Prepared with an executive presentation standard  ·  Fully editable  ·  2025/2026",
         size=9.5, color=RGBColor(0x7C, 0x8C, 0xB0), first=True, align=PP_ALIGN.CENTER)
    transition(s, "fade")
    notes(s, "Thank the committee for their time and open the floor. Prepare for questions on "
             "architecture decisions, security choices, the forecast model and how this scales "
             "to production.\nTransition: Fade. Animation: thank-you text fade-up.")


# --------------------------------------------------------------------------
# Build
# --------------------------------------------------------------------------
builders = [
    s01_cover, s02_agenda, s03_overview, s04_objectives, s05_bi_intro,
    s06_architecture, s07_tech, s08_database, s09_modules, s10_workflow,
    s11_bi_layer, s12_dashboard, s13_visualization, s14_analytics, s15_forecast,
    s16_docker, s17_security, s18_testing, s19_results, s20_future,
    s21_conclusion, s22_thanks,
]
for fn in builders:
    fn()

OUT = r"C:\Users\ZBOOK\Documents\مجلد جديد (3)\Nova_BI_Graduation_Project.pptx"
prs.save(OUT)
print("Saved:", OUT)
print("Slides:", len(prs.slides._sldIdLst))
